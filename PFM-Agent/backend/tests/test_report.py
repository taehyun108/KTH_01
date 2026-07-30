"""
STEP 8 검증: Report Generator 테스트.

  - Word(.docx) / PowerPoint(.pptx) 실제 파일 생성
  - 각주(출처) 포함
  - 개인정보 마스킹 (최종 산출물 방어선)
  - 파일명 보안 (경로 탈출 차단)
  - 소듐이온배터리 시나리오
"""

from __future__ import annotations

from pathlib import Path

import pytest

from app.agents.generator import split_into_sections
from app.mcp_servers.base_server import MCPToolError
from app.mcp_servers.report_server import (
    OUTPUT_DIR,
    ReportMCPServer,
    normalize_citations,
    normalize_sections,
    normalize_slides,
    resolve_output_path,
    safe_filename,
)
from app.scenarios.na_ion_battery import (
    NA_ION_BATTERY,
    SAMPLE_DOCUMENTS,
    create_sample_documents,
    get_scenario,
    list_scenarios,
)


@pytest.fixture(autouse=True)
def clean_output():
    """테스트가 만든 산출물을 정리한다."""
    before = set(OUTPUT_DIR.glob("*")) if OUTPUT_DIR.is_dir() else set()
    yield
    if OUTPUT_DIR.is_dir():
        for path in OUTPUT_DIR.glob("*"):
            if path not in before and path.is_file():
                path.unlink(missing_ok=True)


# ============================================================
# 1. 파일명 보안 (순수 로직)
# ============================================================


def test_safe_filename_strips_special_chars() -> None:
    """특수문자가 제거되어야 한다."""
    name = safe_filename("보고서: 소듐/이온*배터리?", ".docx")
    assert name.endswith(".docx")
    for char in "/*?:":
        assert char not in name


def test_safe_filename_handles_empty_title() -> None:
    """제목이 비어도 기본 이름이 생성되어야 한다."""
    assert safe_filename("", ".docx").endswith(".docx")
    assert "보고서" in safe_filename("!!!", ".docx")


def test_resolve_output_path_blocks_traversal() -> None:
    """
    🔴 보안: 파일명에 경로가 포함돼도 output 폴더를 벗어나면 안 된다.
    """
    path = resolve_output_path("../../../etc/evil.docx", "제목", ".docx")

    assert path.parent == OUTPUT_DIR
    assert path.is_relative_to(OUTPUT_DIR)
    assert ".." not in str(path)


def test_resolve_output_path_adds_extension() -> None:
    """확장자가 없으면 붙여야 한다."""
    path = resolve_output_path("보고서", "제목", ".docx")
    assert path.name.endswith(".docx")


def test_resolve_output_path_auto_names() -> None:
    """파일명을 생략하면 자동 생성되어야 한다."""
    path = resolve_output_path(None, "소듐이온배터리 보고서", ".docx")
    assert path.parent == OUTPUT_DIR
    assert "소듐이온배터리" in path.name


# ============================================================
# 2. 입력 정규화 (순수 로직)
# ============================================================


def test_normalize_sections_masks_personal_info() -> None:
    """🔴 섹션 본문의 개인정보가 마스킹되어야 한다."""
    sections = normalize_sections(
        [{"heading": "담당자", "body": "연락처 010-1234-5678 입니다."}]
    )
    assert "010-1234-5678" not in sections[0]["body"]
    assert "010-****-5678" in sections[0]["body"]


def test_normalize_sections_skips_empty() -> None:
    """빈 섹션은 제외되어야 한다."""
    sections = normalize_sections(
        [{"heading": "", "body": ""}, {"heading": "개요", "body": "내용"}]
    )
    assert len(sections) == 1


def test_normalize_sections_rejects_bad_type() -> None:
    """잘못된 형식은 한글 안내와 함께 거부되어야 한다."""
    with pytest.raises(MCPToolError) as excinfo:
        normalize_sections(["문자열"])  # type: ignore[list-item]
    assert "형식이 잘못" in str(excinfo.value)


def test_normalize_slides_masks_bullets() -> None:
    """슬라이드 글머리 기호도 마스킹되어야 한다."""
    slides = normalize_slides(
        [{"heading": "연락처", "bullets": ["담당자 hong@posco.com", "정상 항목"]}]
    )
    assert "hong@posco.com" not in slides[0]["bullets"][0]
    assert slides[0]["bullets"][1] == "정상 항목"


def test_normalize_citations_deduplicates() -> None:
    """같은 출처는 한 번만 남아야 한다."""
    citations = normalize_citations(
        [
            {"label": "A", "source": "data/policy.md"},
            {"label": "B", "source": "data/policy.md"},
            {"label": "C", "source": "data/market.md"},
        ]
    )
    assert len(citations) == 2


def test_normalize_citations_handles_none() -> None:
    """출처가 없어도 안전해야 한다."""
    assert normalize_citations(None) == []


# ============================================================
# 3. Word 보고서 실제 생성
# ============================================================


async def test_create_word_report_makes_real_file() -> None:
    """Word 보고서 파일이 실제로 만들어져야 한다."""
    server = ReportMCPServer()
    result = await server.handle_call(
        "create_word_report",
        {
            "title": "소듐이온배터리 정책동향 보고서",
            "sections": [
                {"heading": "개요", "body": "소듐이온배터리는 리튬 대비 원가가 낮다."},
                {"heading": "주요 내용", "body": "ESS 분야에서 주목받고 있다."},
                {"heading": "시사점", "body": "선제적 기술 확보가 필요하다."},
            ],
            "citations": [
                {"label": "기술개요", "source": "data/samples/기술개요.md", "quote": "원가가 낮다"},
            ],
        },
    )

    assert result["format"] == "docx"
    assert result["sections"] == 3
    assert result["citations"] == 1

    from app.mcp_servers.base_server import PROJECT_ROOT

    path = PROJECT_ROOT / result["path"]
    assert path.is_file(), "Word 파일이 생성되지 않았습니다"
    assert path.stat().st_size > 0

    # 실제로 Word 문서로 열리는지 확인한다.
    from docx import Document

    document = Document(path)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "소듐이온배터리" in text
    assert "개요" in text
    assert "출처" in text  # 각주 절
    assert "PFM-Agent" in text  # 자동 생성 안내


async def test_word_report_masks_personal_info() -> None:
    """🔴 최종 산출물에서도 개인정보가 가려져야 한다. (마지막 방어선)"""
    server = ReportMCPServer()
    result = await server.handle_call(
        "create_word_report",
        {
            "title": "담당자 명단",
            "sections": [
                {"heading": "연락처", "body": "홍길동 010-1234-5678 hong@posco.com"}
            ],
        },
    )

    assert result["masked_count"] >= 2
    assert "개인정보" in result["message"]

    from docx import Document

    from app.mcp_servers.base_server import PROJECT_ROOT

    document = Document(PROJECT_ROOT / result["path"])
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "010-1234-5678" not in text
    assert "hong@posco.com" not in text


async def test_word_report_requires_sections() -> None:
    """본문이 없으면 한글 안내와 함께 거부되어야 한다."""
    server = ReportMCPServer()
    with pytest.raises(MCPToolError) as excinfo:
        await server.handle_call(
            "create_word_report", {"title": "제목", "sections": []}
        )
    assert "본문이 없습니다" in str(excinfo.value)


# ============================================================
# 4. PowerPoint 실제 생성
# ============================================================


async def test_create_ppt_report_makes_real_file() -> None:
    """PPT 파일이 실제로 만들어져야 한다."""
    server = ReportMCPServer()
    result = await server.handle_call(
        "create_ppt_report",
        {
            "title": "소듐이온배터리 동향",
            "slides": [
                {"heading": "기술 개요", "bullets": ["원가가 낮다", "ESS 에 적합"]},
                {"heading": "시사점", "bullets": ["선제 기술 확보 필요"]},
            ],
            "citations": [{"label": "시장동향", "source": "data/market.md"}],
        },
    )

    assert result["format"] == "pptx"
    # 표지 + 본문 2 + 출처 1 = 4
    assert result["slides"] == 4

    from pptx import Presentation

    from app.mcp_servers.base_server import PROJECT_ROOT

    path = PROJECT_ROOT / result["path"]
    assert path.is_file(), "PPT 파일이 생성되지 않았습니다"

    presentation = Presentation(path)
    assert len(presentation.slides) == 4

    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    combined = "\n".join(texts)

    assert "소듐이온배터리" in combined
    assert "기술 개요" in combined
    assert "출처" in combined


async def test_ppt_requires_slides() -> None:
    """슬라이드가 없으면 거부되어야 한다."""
    server = ReportMCPServer()
    with pytest.raises(MCPToolError) as excinfo:
        await server.handle_call("create_ppt_report", {"title": "제목", "slides": []})
    assert "슬라이드가 없습니다" in str(excinfo.value)


async def test_unknown_report_tool_rejected() -> None:
    """없는 도구는 거부되어야 한다."""
    server = ReportMCPServer()
    with pytest.raises(MCPToolError) as excinfo:
        await server.handle_call("create_pdf_report", {})
    assert "도구가 없습니다" in str(excinfo.value)


# ============================================================
# 5. Generator 본문 → 섹션 변환
# ============================================================


def test_split_into_sections_by_heading() -> None:
    """마크다운 소제목 기준으로 섹션이 나뉘어야 한다."""
    body = "## 개요\n\n첫 내용\n\n## 시사점\n\n둘째 내용"
    sections = split_into_sections(body)

    assert len(sections) == 2
    assert sections[0]["heading"] == "개요"
    assert "첫 내용" in sections[0]["body"]
    assert sections[1]["heading"] == "시사점"


def test_split_into_sections_without_heading() -> None:
    """소제목이 없으면 통째로 '본문' 하나가 되어야 한다."""
    sections = split_into_sections("소제목 없는 그냥 긴 글입니다.")
    assert len(sections) == 1
    assert sections[0]["heading"] == "본문"


def test_split_into_sections_empty() -> None:
    """빈 본문도 안전하게 처리되어야 한다."""
    sections = split_into_sections("")
    assert len(sections) == 1


# ============================================================
# 6. 시나리오
# ============================================================


def test_scenario_registered() -> None:
    """소듐이온배터리 시나리오가 등록되어 있어야 한다."""
    scenario = get_scenario("na_ion_battery")
    assert scenario is not None
    assert scenario.key == NA_ION_BATTERY.key
    assert "소듐이온배터리" in scenario.request
    assert "출처" in scenario.request  # 출처 표기를 요구


def test_list_scenarios() -> None:
    """시나리오 목록을 API 형식으로 반환할 수 있어야 한다."""
    scenarios = list_scenarios()
    assert len(scenarios) >= 1
    assert scenarios[0]["key"] == "na_ion_battery"


def test_sample_documents_are_marked_as_examples() -> None:
    """
    🔴 샘플 자료는 실제 정책 자료로 오인되면 안 되므로
    모든 문서에 '예시' 경고가 있어야 한다.
    """
    for name, content in SAMPLE_DOCUMENTS.items():
        assert "예시" in content, f"{name} 에 예시 표시가 없습니다"
        assert "실제" in content, f"{name} 에 주의 문구가 없습니다"


def test_create_sample_documents(tmp_path: Path) -> None:
    """샘플 자료 파일이 생성되어야 한다."""
    created = create_sample_documents(tmp_path)

    assert len(created) == len(SAMPLE_DOCUMENTS)
    for path in created:
        assert path.is_file()
        assert "소듐이온배터리" in path.read_text(encoding="utf-8") or "이차전지" in path.read_text(
            encoding="utf-8"
        )


# ============================================================
# 7. 보고서 제목 정리
# ============================================================


@pytest.mark.parametrize(
    ("request_text", "expected"),
    [
        (
            "소듐이온배터리 관련 국내외 정책 동향을 조사해서 보고서를 작성해줘. 출처도 표기해줘.",
            "소듐이온배터리 관련 국내외 정책 동향",
        ),
        ("소듐이온배터리 정책동향 보고서를 만들어줘", "소듐이온배터리 정책동향 보고서"),
        ("ESS 시장 현황 정리해줘", "ESS 시장 현황"),
    ],
)
def test_clean_title_removes_instructions(request_text: str, expected: str) -> None:
    """
    사용자 요청문을 그대로 제목에 넣으면 '~해줘' 같은 구어체가 남아
    정부 부처 대응 문서로 어색하다. 지시 표현이 제거되어야 한다.
    """
    from app.agents.generator import clean_title

    assert clean_title(request_text) == expected


def test_clean_title_handles_empty() -> None:
    """빈 요청도 안전하게 처리되어야 한다."""
    from app.agents.generator import clean_title

    assert clean_title("") == "자동 생성 보고서"
    assert clean_title("   ") == "자동 생성 보고서"


def test_clean_title_length_capped() -> None:
    """제목이 지나치게 길어지지 않아야 한다."""
    from app.agents.generator import clean_title

    assert len(clean_title("가" * 200)) <= 60
