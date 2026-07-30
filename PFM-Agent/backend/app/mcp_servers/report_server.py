"""
report MCP 서버 — 보고서 생성 도구. (STEP 8 구현)

제공 도구:
  - create_word_report(title, sections, citations) : Word(.docx) 생성
  - create_ppt_report(title, slides, citations)    : PowerPoint(.pptx) 생성

산출물은 모두 프로젝트 폴더 내부(`./output`)에 저장한다. (포터블 원칙)

각주(인용) 처리
---------------
근거를 사용한 보고서는 반드시 출처를 남긴다.
본문 뒤에 "출처" 절을 두고 번호를 매겨 정리한다.
정부 부처 대응 문서는 근거 추적이 중요하기 때문이다.

개인정보 보호
-------------
보고서에 들어가는 모든 텍스트는 **마스킹을 거친다.**
Verifier 단계에서 이미 걸러지지만, 최종 산출물에서 한 번 더 확인한다.
(마지막 방어선 — 보고서는 외부로 나가는 문서이므로)

단독 실행:
    python -m app.mcp_servers.report_server
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

from mcp.types import Tool

from app.mcp_servers.base_server import (
    PROJECT_ROOT,
    BaseMCPServer,
    MCPToolError,
    to_relative,
)
from app.security.masking import mask_text

#: 산출물 저장 폴더
OUTPUT_DIR: Path = PROJECT_ROOT / "output"

#: 파일명 최대 길이
MAX_FILENAME_LENGTH: int = 60

#: 인용(각주) 항목 스키마 — Word/PPT 공통
_CITATION_SCHEMA: dict[str, Any] = {
    "type": "array",
    "description": "각주로 넣을 출처 목록. 근거가 있는 내용에는 반드시 출처를 붙이세요.",
    "items": {
        "type": "object",
        "properties": {
            "label": {"type": "string", "description": "출처 표시 이름"},
            "source": {"type": "string", "description": "출처 경로 또는 URL"},
            "quote": {"type": "string", "description": "인용한 원문 (선택)"},
        },
        "required": ["label", "source"],
    },
}


# ============================================================
# 순수 로직 (환경 무관 — 단독 테스트 가능)
# ============================================================


def safe_filename(title: str, extension: str) -> str:
    """
    제목을 안전한 파일명으로 바꾼다.

    경로 구분자·특수문자를 제거해 폴더 탈출을 막는다.
    """
    cleaned = re.sub(r"[^\w가-힣 _-]", "", title).strip()
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = cleaned[:MAX_FILENAME_LENGTH] or "보고서"

    timestamp = datetime.now().astimezone().strftime("%Y%m%d_%H%M%S")
    suffix = extension if extension.startswith(".") else f".{extension}"
    return f"{timestamp}_{cleaned}{suffix}"


def resolve_output_path(filename: str | None, title: str, extension: str) -> Path:
    """
    저장 경로를 결정한다.

    사용자가 파일명을 지정해도 **경로 부분은 무시**하고 파일명만 사용한다.
    (`../../windows/system32/evil.docx` 같은 시도를 막기 위함)
    """
    if filename and filename.strip():
        # 경로 구분자를 제거해 파일명만 취한다.
        base = Path(filename.strip()).name
        base = re.sub(r"[^\w가-힣 ._-]", "", base).strip()
        if not base:
            base = safe_filename(title, extension)
        elif not base.lower().endswith(extension.lower()):
            base = f"{base}{extension}"
    else:
        base = safe_filename(title, extension)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR / base


def normalize_sections(sections: list[dict[str, Any]]) -> list[dict[str, str]]:
    """
    본문 섹션을 검증하고 마스킹한다.

    Raises:
        MCPToolError: 형식이 잘못된 경우
    """
    if not isinstance(sections, list):
        raise MCPToolError("sections 는 목록이어야 합니다.")

    normalized: list[dict[str, str]] = []
    for index, section in enumerate(sections, start=1):
        if not isinstance(section, dict):
            raise MCPToolError(f"{index}번째 섹션의 형식이 잘못되었습니다.")

        heading = str(section.get("heading", "")).strip()
        body = str(section.get("body", "")).strip()
        if not heading and not body:
            continue

        normalized.append(
            {
                "heading": mask_text(heading).text,
                "body": mask_text(body).text,
            }
        )

    return normalized


def normalize_slides(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """
    슬라이드를 검증하고 마스킹한다.

    Raises:
        MCPToolError: 형식이 잘못된 경우
    """
    if not isinstance(slides, list):
        raise MCPToolError("slides 는 목록이어야 합니다.")

    normalized: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            raise MCPToolError(f"{index}번째 슬라이드의 형식이 잘못되었습니다.")

        heading = str(slide.get("heading", "")).strip()
        raw_bullets = slide.get("bullets", [])
        if not isinstance(raw_bullets, list):
            raise MCPToolError(f"{index}번째 슬라이드의 bullets 는 목록이어야 합니다.")

        bullets = [
            mask_text(str(bullet).strip()).text
            for bullet in raw_bullets
            if str(bullet).strip()
        ]
        if not heading and not bullets:
            continue

        normalized.append({"heading": mask_text(heading).text, "bullets": bullets})

    return normalized


def normalize_citations(citations: list[dict[str, Any]] | None) -> list[dict[str, str]]:
    """출처 목록을 정리한다. (중복 제거 + 마스킹)"""
    if not citations:
        return []

    seen: set[str] = set()
    normalized: list[dict[str, str]] = []

    for citation in citations:
        if not isinstance(citation, dict):
            continue
        source = str(citation.get("source", "")).strip()
        if not source or source in seen:
            continue
        seen.add(source)

        normalized.append(
            {
                "label": mask_text(str(citation.get("label", source)).strip()).text,
                "source": mask_text(source).text,
                "quote": mask_text(str(citation.get("quote", "")).strip()).text,
            }
        )

    return normalized


def count_personal_info(sections: list[dict[str, Any]]) -> int:
    """섹션 전체에서 마스킹된 개인정보 건수를 센다. (사용자 안내용)"""
    total = 0
    for section in sections:
        for value in section.values():
            if isinstance(value, str):
                total += mask_text(value).total
            elif isinstance(value, list):
                for item in value:
                    total += mask_text(str(item)).total
    return total


# ============================================================
# MCP 서버
# ============================================================


class ReportMCPServer(BaseMCPServer):
    """Word / PowerPoint 보고서 생성 도구를 제공하는 MCP 서버."""

    def __init__(self) -> None:
        super().__init__("report")

    def get_tools(self) -> list[Tool]:
        """이 서버가 제공하는 도구 목록."""
        return [
            Tool(
                name="create_word_report",
                description=(
                    "제목과 본문 섹션으로 Word(.docx) 보고서를 만들어 "
                    "./output 폴더에 저장하고 파일 경로를 반환합니다. "
                    "정부 부처 대응 보고서 등 문서 형태의 최종 산출물을 만들 때 사용하세요."
                ),
                inputSchema={
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "보고서 제목"},
                        "sections": {
                            "type": "array",
                            "description": "본문 섹션 목록 (순서대로 문서에 작성됨)",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "heading": {"type": "string", "description": "섹션 제목"},
                                    "body": {"type": "string", "description": "섹션 본문"},
                                },
                                "required": ["heading", "body"],
                            },
                        },
                        "citations": _CITATION_SCHEMA,
                        "filename": {
                            "type": "string",
                            "description": "저장할 파일 이름 (생략 시 제목과 날짜로 자동 생성)",
                        },
                    },
                    "required": ["title", "sections"],
                },
            ),
            Tool(
                name="create_ppt_report",
                description=(
                    "제목과 슬라이드 목록으로 PowerPoint(.pptx) 보고서를 만들어 "
                    "./output 폴더에 저장하고 파일 경로를 반환합니다. "
                    "발표용 요약 자료가 필요할 때 사용하세요."
                ),
                inputSchema={
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "발표 자료 제목"},
                        "slides": {
                            "type": "array",
                            "description": "슬라이드 목록 (순서대로 생성됨)",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "heading": {"type": "string", "description": "슬라이드 제목"},
                                    "bullets": {
                                        "type": "array",
                                        "description": "본문 글머리 기호 목록",
                                        "items": {"type": "string"},
                                    },
                                },
                                "required": ["heading", "bullets"],
                            },
                        },
                        "citations": _CITATION_SCHEMA,
                        "filename": {
                            "type": "string",
                            "description": "저장할 파일 이름 (생략 시 제목과 날짜로 자동 생성)",
                        },
                    },
                    "required": ["title", "slides"],
                },
            ),
        ]

    async def handle_call(self, name: str, arguments: dict[str, Any]) -> Any:
        """도구 호출을 처리한다."""
        if name == "create_word_report":
            return self._create_word(
                title=str(arguments["title"]),
                sections=arguments["sections"],
                citations=arguments.get("citations"),
                filename=arguments.get("filename"),
            )
        if name == "create_ppt_report":
            return self._create_ppt(
                title=str(arguments["title"]),
                slides=arguments["slides"],
                citations=arguments.get("citations"),
                filename=arguments.get("filename"),
            )
        raise MCPToolError(f"report 서버에 '{name}' 도구가 없습니다.")

    # ------------------------------------------------------------
    # Word 보고서
    # ------------------------------------------------------------
    def _create_word(
        self,
        *,
        title: str,
        sections: list[dict[str, Any]],
        citations: list[dict[str, Any]] | None,
        filename: str | None,
    ) -> dict[str, Any]:
        """Word(.docx) 보고서를 생성한다."""
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Pt
        except ImportError as exc:
            raise MCPToolError(
                "Word 문서 생성 라이브러리(python-docx)가 설치되어 있지 않습니다.\n"
                "  해결: first_time_setup.bat 를 다시 실행하세요."
            ) from exc

        clean_title = mask_text(title.strip() or "보고서").text
        clean_sections = normalize_sections(sections)
        clean_citations = normalize_citations(citations)
        masked_count = count_personal_info(sections) + count_personal_info(
            citations or []
        )

        if not clean_sections:
            raise MCPToolError(
                "보고서에 넣을 본문이 없습니다. sections 에 최소 1개 항목이 필요합니다."
            )

        document = Document()

        # --- 제목 ---
        heading = document.add_heading(clean_title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # --- 작성 정보 ---
        info = document.add_paragraph()
        info.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        info_run = info.add_run(
            f"작성일: {datetime.now().astimezone().strftime('%Y년 %m월 %d일')}"
        )
        info_run.font.size = Pt(9)

        # --- 본문 ---
        for section in clean_sections:
            if section["heading"]:
                document.add_heading(section["heading"], level=1)
            if section["body"]:
                for paragraph_text in section["body"].split("\n"):
                    text = paragraph_text.strip()
                    if text:
                        document.add_paragraph(text)

        # --- 출처 (각주) ---
        if clean_citations:
            document.add_page_break()
            document.add_heading("출처", level=1)
            for index, citation in enumerate(clean_citations, start=1):
                paragraph = document.add_paragraph()
                paragraph.add_run(f"[{index}] {citation['label']}").bold = True
                paragraph.add_run(f"\n     {citation['source']}")
                if citation["quote"]:
                    quote = document.add_paragraph(f'     "{citation["quote"]}"')
                    quote.runs[0].italic = True

        # --- 안내 문구 ---
        document.add_paragraph()
        notice = document.add_paragraph()
        notice_run = notice.add_run(
            "※ 본 보고서는 PFM-Agent 가 자동 생성했습니다. 최종 검토 책임은 작성자에게 있습니다."
        )
        notice_run.font.size = Pt(9)
        notice_run.italic = True

        path = resolve_output_path(filename, clean_title, ".docx")
        try:
            document.save(path)
        except OSError as exc:
            raise MCPToolError(f"보고서를 저장하지 못했습니다: {exc}") from exc

        message = f"Word 보고서를 생성했습니다: {to_relative(path)}"
        if masked_count:
            message += f" (개인정보 {masked_count}건은 가려서 저장했습니다)"

        return {
            "path": to_relative(path),
            "format": "docx",
            "sections": len(clean_sections),
            "citations": len(clean_citations),
            "masked_count": masked_count,
            "message": message,
        }

    # ------------------------------------------------------------
    # PowerPoint 보고서
    # ------------------------------------------------------------
    def _create_ppt(
        self,
        *,
        title: str,
        slides: list[dict[str, Any]],
        citations: list[dict[str, Any]] | None,
        filename: str | None,
    ) -> dict[str, Any]:
        """PowerPoint(.pptx) 보고서를 생성한다."""
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError as exc:
            raise MCPToolError(
                "PPT 생성 라이브러리(python-pptx)가 설치되어 있지 않습니다.\n"
                "  해결: first_time_setup.bat 를 다시 실행하세요."
            ) from exc

        clean_title = mask_text(title.strip() or "발표 자료").text
        clean_slides = normalize_slides(slides)
        clean_citations = normalize_citations(citations)
        masked_count = count_personal_info(slides) + count_personal_info(citations or [])

        if not clean_slides:
            raise MCPToolError(
                "발표 자료에 넣을 슬라이드가 없습니다. slides 에 최소 1개 항목이 필요합니다."
            )

        presentation = Presentation()

        # --- 표지 ---
        cover_layout = presentation.slide_layouts[0]
        cover = presentation.slides.add_slide(cover_layout)
        cover.shapes.title.text = clean_title
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = (
                f"작성일: {datetime.now().astimezone().strftime('%Y년 %m월 %d일')}"
            )

        # --- 본문 슬라이드 ---
        content_layout = presentation.slide_layouts[1]
        for slide_data in clean_slides:
            slide = presentation.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data["heading"]

            body = slide.placeholders[1].text_frame
            body.clear()
            for index, bullet in enumerate(slide_data["bullets"]):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(18)

        # --- 출처 슬라이드 ---
        if clean_citations:
            source_slide = presentation.slides.add_slide(content_layout)
            source_slide.shapes.title.text = "출처"
            body = source_slide.placeholders[1].text_frame
            body.clear()
            for index, citation in enumerate(clean_citations, start=1):
                paragraph = body.paragraphs[0] if index == 1 else body.add_paragraph()
                paragraph.text = f"[{index}] {citation['label']} — {citation['source']}"
                paragraph.font.size = Pt(14)

        path = resolve_output_path(filename, clean_title, ".pptx")
        try:
            presentation.save(path)
        except OSError as exc:
            raise MCPToolError(f"발표 자료를 저장하지 못했습니다: {exc}") from exc

        message = f"PPT 발표 자료를 생성했습니다: {to_relative(path)}"
        if masked_count:
            message += f" (개인정보 {masked_count}건은 가려서 저장했습니다)"

        return {
            "path": to_relative(path),
            "format": "pptx",
            "slides": len(clean_slides) + 1 + (1 if clean_citations else 0),
            "citations": len(clean_citations),
            "masked_count": masked_count,
            "message": message,
        }


def main() -> None:
    """단독 실행 진입점."""
    ReportMCPServer.main()


if __name__ == "__main__":
    main()
